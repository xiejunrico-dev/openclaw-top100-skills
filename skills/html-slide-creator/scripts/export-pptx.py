#!/usr/bin/env python3
"""
export-pptx.py — Export slide-creator HTML presentations to PPTX.

Uses Playwright with your existing system Chrome to take pixel-perfect
screenshots of each slide, then assembles them into a PowerPoint file.
No Chromium download needed if Chrome/Edge/Brave is already installed.

Usage:
    python export-pptx.py <presentation.html> [output.pptx] [--width W] [--height H]

Dependencies:
    pip install playwright python-pptx
    (No browser download needed if Chrome is already installed)
"""

import sys
import os
import argparse
import tempfile
import shutil
from pathlib import Path


# ─── Dependency check ─────────────────────────────────────────────────────────

def check_deps():
    missing = []
    try:
        from playwright.sync_api import sync_playwright  # noqa
    except ImportError:
        missing.append("playwright")
    try:
        from pptx import Presentation  # noqa
    except ImportError:
        missing.append("python-pptx")
    if missing:
        print(f"Missing dependencies. Install with:")
        print(f"  pip install {' '.join(missing)}")
        sys.exit(1)

check_deps()

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches


# ─── Screenshot via system browser ────────────────────────────────────────────

# CSS injected before screenshotting to:
# 1. Disable all animations/transitions so content is immediately fully visible
# 2. Force .reveal elements (which start invisible) to be shown
DISABLE_ANIMATIONS_CSS = """
*, *::before, *::after {
    animation-duration: 0.01ms !important;
    animation-delay: 0.01ms !important;
    transition-duration: 0.01ms !important;
    transition-delay: 0.01ms !important;
}
.reveal, [class*="reveal"] {
    opacity: 1 !important;
    transform: none !important;
    filter: none !important;
}
"""


def find_and_launch_browser(playwright):
    """
    Try channels in order of preference: use system-installed browsers first
    to avoid downloading anything. Falls back to Playwright's own Chromium
    only as a last resort.

    Linux servers: add --no-sandbox because root/container environments
    block Chrome's sandboxing by default.
    """
    import platform
    is_linux = platform.system() == 'Linux'
    # --no-sandbox is required in Docker/root/CI environments on Linux
    extra_args = ['--no-sandbox', '--disable-setuid-sandbox'] if is_linux else []

    channels = ['chrome', 'msedge', 'chromium']
    for channel in channels:
        try:
            browser = playwright.chromium.launch(
                channel=channel, headless=True, args=extra_args
            )
            print(f"  Using browser: {channel}")
            return browser
        except Exception:
            continue

    # Last resort: Playwright's own Chromium (may require: playwright install chromium)
    try:
        browser = playwright.chromium.launch(headless=True, args=extra_args)
        print("  Using browser: playwright-chromium")
        return browser
    except Exception as e:
        print("\nNo browser found.")
        if is_linux:
            print("  Linux options (pick one):")
            print("    apt install chromium-browser      # system package (~100MB)")
            print("    playwright install chromium        # self-contained (~170MB)")
        else:
            print("  Install Google Chrome: https://www.google.com/chrome/")
            print("  Or run: playwright install chromium")
        raise SystemExit(1) from e


def screenshot_slides(html_path, output_dir, width=1440, height=900):
    """
    Open the HTML presentation in a headless browser, scroll to each .slide,
    and save a screenshot. Returns a list of screenshot paths in slide order.
    """
    html_path = Path(html_path).resolve()
    url = f"file://{html_path}"
    screenshots = []

    with sync_playwright() as p:
        browser = find_and_launch_browser(p)
        page = browser.new_page(viewport={"width": width, "height": height})

        # Navigate and wait for fonts + images
        page.goto(url, wait_until="networkidle")
        page.wait_for_timeout(800)

        # Inject CSS to disable animations so all content is immediately visible
        page.add_style_tag(content=DISABLE_ANIMATIONS_CSS)

        # Trigger the Intersection Observer by marking all slides visible
        # (slide-creator adds .visible when a slide enters the viewport)
        page.evaluate("""
            document.querySelectorAll('.slide').forEach(s => s.classList.add('visible'));
        """)
        page.wait_for_timeout(200)

        slide_count = page.evaluate("document.querySelectorAll('.slide').length")
        if slide_count == 0:
            print("No .slide elements found in HTML.")
            browser.close()
            return []

        print(f"  Found {slide_count} slides. Capturing...")

        for i in range(slide_count):
            # Scroll this slide into view
            page.evaluate(f"document.querySelectorAll('.slide')[{i}].scrollIntoView()")
            page.wait_for_timeout(150)

            # Get title for progress display
            title = page.evaluate(f"""
                (() => {{
                    const s = document.querySelectorAll('.slide')[{i}];
                    const h = s.querySelector('h1,h2,h3');
                    return h ? h.textContent.trim().slice(0, 50) : 'Slide {i+1}';
                }})()
            """)
            print(f"  [{i+1}/{slide_count}] {title}")

            img_path = output_dir / f"slide_{i:03d}.png"
            page.screenshot(path=str(img_path), full_page=False)
            screenshots.append(img_path)

        browser.close()

    return screenshots


# ─── Assemble screenshots into PPTX ──────────────────────────────────────────

def assemble_pptx(screenshots, output_path, width=1440, height=900):
    """
    Create a PPTX where each slide is the full-bleed screenshot.
    Aspect ratio matches the screenshot dimensions.
    """
    # Scale to PowerPoint inches maintaining 16:9 (or actual ratio)
    ratio = height / width
    slide_w = Inches(13.33)
    slide_h = Inches(13.33 * ratio)

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h

    blank_layout = prs.slide_layouts[6]  # blank

    for img_path in screenshots:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(img_path), 0, 0, slide_w, slide_h)

    prs.save(str(output_path))


# ─── Main ─────────────────────────────────────────────────────────────────────

def export(html_path, output_path=None, width=1440, height=900):
    html_path = Path(html_path).resolve()
    if not html_path.exists():
        print(f"Error: file not found: {html_path}")
        sys.exit(1)

    output_path = Path(output_path) if output_path else html_path.with_suffix('.pptx')

    print(f"Exporting: {html_path.name}")
    print(f"Viewport:  {width}×{height}")

    tmp_dir = Path(tempfile.mkdtemp(prefix="slide-export-"))
    try:
        screenshots = screenshot_slides(html_path, tmp_dir, width, height)
        if not screenshots:
            print("Nothing to export.")
            return

        print(f"\nAssembling PPTX...")
        assemble_pptx(screenshots, output_path, width, height)
        print(f"✓ Saved: {output_path}  ({len(screenshots)} slides)")
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

    return output_path


def main():
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("html", help="Path to the HTML presentation")
    parser.add_argument("output", nargs="?", help="Output .pptx path (default: same name as HTML)")
    parser.add_argument("--width",  type=int, default=1440, help="Viewport width  (default: 1440)")
    parser.add_argument("--height", type=int, default=900,  help="Viewport height (default: 900)")
    args = parser.parse_args()
    export(args.html, args.output, args.width, args.height)


if __name__ == "__main__":
    main()
